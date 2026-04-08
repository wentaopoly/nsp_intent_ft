"""
Generate the NSP Intent Fine-Tuning presentation (PPTX).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# --- Color scheme ---
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x00, 0x7B, 0xFF)
GREEN = RGBColor(0x28, 0xA7, 0x45)
RED = RGBColor(0xDC, 0x35, 0x45)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
DARK_TEXT = RGBColor(0x21, 0x25, 0x29)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xAD, 0xB5, 0xBD)
    p2.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Bullets
    body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.5), Inches(5))
    tf2 = body_box.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        if isinstance(bullet, tuple):
            text, level = bullet
        else:
            text, level = bullet, 0

        p.text = text
        p.font.size = Pt(16) if level == 0 else Pt(14)
        p.font.color.rgb = DARK_TEXT if level == 0 else GRAY
        p.space_after = Pt(6)
        p.level = level

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(8.5), Inches(0.5))
        tf3 = note_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(11)
        p3.font.italic = True
        p3.font.color.rgb = GRAY

    return slide


def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(4.5), Inches(5.5))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = left_title
    p_l.font.size = Pt(18)
    p_l.font.bold = True
    p_l.font.color.rgb = RED
    for item in left_items:
        p = tf_l.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(4)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(5.5))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = right_title
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = GREEN
    for item in right_items:
        p = tf_r.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(4)

    return slide


def add_code_slide(prs, title, code_text, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.4))
        tfs = sub_box.text_frame
        ps = tfs.paragraphs[0]
        ps.text = subtitle
        ps.font.size = Pt(14)
        ps.font.color.rgb = GRAY

    top = 1.2 if subtitle else 0.9
    code_box = slide.shapes.add_textbox(Inches(0.4), Inches(top), Inches(9.2), Inches(6.0 - top))
    tf2 = code_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = code_text
    p2.font.size = Pt(10)
    p2.font.name = "Consolas"
    p2.font.color.rgb = DARK_TEXT

    return slide


def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.4 * n_rows))
    table = table_shape.table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = DARK_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === Slide 1: Title ===
    add_title_slide(
        prs,
        "NSP Intent JSON Generation",
        "Fine-Tuning Qwen3.5-9B for Automated Network Service Configuration\n"
        "Nokia NAR Lab | transport-intent-transformer"
    )

    # === Slide 2: Problem Statement ===
    add_content_slide(prs, "The Problem", [
        "Nokia NSP requires precise JSON intent configurations to deploy network services",
        "Engineers must manually write 30-900 lines of JSON for each service deployment",
        "Each field (IP addresses, ports, VLANs, SDPs) must be exactly correct",
        "One mistake = deployment failure",
        "",
        "Goal: Let engineers describe services in natural language,",
        "and automatically generate correct, API-ready JSON configurations",
    ])

    # === Slide 3: Previous Approach (Marla) ===
    add_section_slide(prs, "Previous Approach & Its Problems")

    add_comparison_slide(
        prs,
        "Marla's Approach vs Our Approach",
        "Marla's Approach (Problems)",
        [
            "Trained on documentation Q&A (95 samples)",
            "  -> Model learned to answer questions, not fill JSON",
            "",
            "Two-step pipeline used base model, not fine-tuned",
            "  -> Fine-tuning had zero impact on inference",
            "",
            "Preprocessing bug lost context information",
            "  -> combined_prompts created but never used",
            "",
            "15 epochs on 76 samples = severe overfitting",
            "  -> Loss 0.09, Accuracy 97.8% (memorized)",
            "",
            "VPRN JSON too large (7400 tokens) for output",
        ],
        "Our Approach (Solutions)",
        [
            "Train directly on target task: NL -> fill-values",
            "  -> Training = Inference, perfectly aligned",
            "",
            "Fill-Values format (compact, ~200 tokens)",
            "  -> Template merging handles boilerplate",
            "",
            "1500 synthetic samples from real constraints",
            "  -> Sarvesh's operational rules built in",
            "",
            "5 epochs + eval monitoring + early stopping",
            "  -> No overfitting (eval_loss < train_loss)",
            "",
            "All intent types fit within 2048 tokens",
        ],
    )

    # === Slide 5: Architecture ===
    add_section_slide(prs, "Architecture: Fill-Values Approach")

    add_content_slide(prs, "The Fill-Values Architecture", [
        "Instead of generating complete JSON (30-900 lines), the model outputs only the values that need to change",
        "",
        "Pipeline:  User Instruction  ->  Model (fill-values)  ->  Template Merge  ->  API-ready JSON",
        "",
        "Advantages:",
        ("Model focuses purely on semantic understanding (NL -> parameters)", 1),
        ("Structural correctness guaranteed by template (not by model)", 1),
        ("All intent types (including VPRN) fit within 2048 tokens", 1),
        ("Deterministic merge script = no hallucinated JSON structure", 1),
    ])

    # === Slide 6: Input/Output Example ===
    add_code_slide(
        prs,
        "Input / Output Example",
        'INPUT (Natural Language):\n'
        '"Create an E-Pipe service named \'Epipe-VLAN-1001-nvlink\' for customer 10.\n'
        ' Connect device 192.168.0.37 port 1/2/c4/1 to device 192.168.0.16\n'
        ' port 1/2/c5/1, VLAN 1001, MTU 1492. SDP 3716 and 1637."\n\n'
        'OUTPUT (Model generates fill-values):\n'
        '{\n'
        '  "intent_type": "epipe",\n'
        '  "template_name": "ePIPE-Service-Using-SDP",\n'
        '  "fill_values": {\n'
        '    "service-name": "Epipe-VLAN-1001-nvlink",\n'
        '    "customer-id": 10,\n'
        '    "site-a.device-id": "192.168.0.37",\n'
        '    "site-a.endpoint[0].port-id": "1/2/c4/1",\n'
        '    "site-a.endpoint[0].outer-vlan-tag": 1001,\n'
        '    "site-b.device-id": "192.168.0.16",\n'
        '    "site-b.endpoint[0].port-id": "1/2/c5/1",\n'
        '    "sdp[0].sdp-id": "3716",\n'
        '    "sdp[0].source-device-id": "192.168.0.37",\n'
        '    "sdp[0].destination-device-id": "192.168.0.16",\n'
        '    "sdp[1].sdp-id": "1637",\n'
        '    "sdp[1].source-device-id": "192.168.0.16",\n'
        '    "sdp[1].destination-device-id": "192.168.0.37"\n'
        '  }\n'
        '}\n\n'
        'FINAL (After template merge -> API-ready JSON, 66 lines, directly POST-able to NSP)',
        subtitle="The model extracts parameters; a deterministic script produces the final JSON"
    )

    # === Slide 7: Training Data ===
    add_section_slide(prs, "Training Data Generation")

    add_content_slide(prs, "Data Sources & Constraints", [
        "Real operational data from Sarvesh's NSP deployments:",
        ("Complete cURL commands with actual JSON payloads (tunnel, epipe, VPRN)", 1),
        ("Real device IPs, port formats, naming conventions", 1),
        ("SDP bidirectionality rule: two entries with swapped src/dst", 1),
        ("SDP ID derivation: last IP octets concatenated (37+16 = 3716)", 1),
        "",
        "We extracted these constraints and built a synthetic data generator:",
        ("30-50 natural language templates per intent type (formal, conversational, terse)", 1),
        ("Random but realistic values within valid ranges", 1),
        ("Automatic validation: IP format, port format, VLAN range, SDP rules", 1),
        ("100% pass rate on all 1500 generated samples", 1),
    ])

    add_table_slide(prs, "Training Data Distribution",
        ["Intent Type", "Count", "Complexity", "Fillable Fields"],
        [
            ["Tunnel", "400", "Simple", "6 fields"],
            ["EPIPE", "600", "Medium", "14 fields"],
            ["VPRN (1 site)", "300", "Medium", "~15 fields"],
            ["VPRN (2 sites)", "200", "Complex", "~30 fields"],
            ["Total", "1,500", "-", "-"],
        ]
    )

    # === Slide 9: Training Configuration ===
    add_section_slide(prs, "Training Configuration & Results")

    add_table_slide(prs, "Training Setup",
        ["Parameter", "Value", "Rationale"],
        [
            ["Base Model", "Qwen3.5-9B", "Strong structured output capability"],
            ["Precision", "BF16 (no quantization)", "2x48GB GPUs = plenty of VRAM"],
            ["Method", "LoRA (r=32, alpha=64)", "0.65% trainable params"],
            ["Target Modules", "q/k/v/o_proj + gate/up/down_proj", "Attention + MLP for JSON output"],
            ["GPUs", "2x NVIDIA RTX 6000 Ada (DDP)", "Each loads full model copy"],
            ["Epochs", "5", "Prevents overfitting"],
            ["Learning Rate", "2e-4 (cosine decay)", "Conservative, stable"],
            ["Effective Batch Size", "32 (2x2x8)", "batch=2, grad_accum=8, 2 GPUs"],
            ["Training Time", "52 minutes", "190 steps total"],
        ]
    )

    # === Slide 11: Training Curves ===
    add_table_slide(prs, "Training Progress",
        ["Step", "Epoch", "Train Loss", "Train Acc", "Eval Loss", "Eval Acc"],
        [
            ["10", "0.27", "0.793", "80.5%", "-", "-"],
            ["20", "0.53", "0.258", "92.3%", "-", "-"],
            ["30", "0.80", "0.163", "94.1%", "-", "-"],
            ["40", "1.05", "0.139", "94.7%", "-", "-"],
            ["50", "1.32", "0.126", "95.1%", "0.120", "95.3%"],
            ["100", "2.64", "0.119", "95.2%", "0.115", "95.4%"],
            ["150", "3.95", "0.115", "95.4%", "0.114", "95.4%"],
            ["190", "5.00", "0.115", "95.4%", "-", "-"],
        ]
    )

    add_content_slide(prs, "Key Training Observations", [
        "Eval Loss consistently lower than Train Loss",
        ("No overfitting - model generalizes well to unseen data", 1),
        "",
        "Fast convergence: 94% accuracy within first epoch",
        ("Fill-values format makes the task learnable and structured", 1),
        "",
        "Stable plateau at 95.4% accuracy from epoch 2 onward",
        ("Model learned all patterns; further epochs provide fine-tuning", 1),
        "",
        "Contrast with Marla: her loss 0.09 / acc 97.8% = overfitting (76 samples, 15 epochs)",
        ("Our loss 0.115 / acc 95.4% = healthy generalization (1200 samples, 5 epochs)", 1),
    ])

    # === Slide 13: Results ===
    add_section_slide(prs, "Results: Real-World Validation")

    add_content_slide(prs, "Validation Against Sarvesh's Real Deployments", [
        "Test 1: MPLS Tunnel",
        ('Input: "Create an MPLS tunnel from 192.168.0.16 to 192.168.0.37, SDP 1637"', 1),
        ("Result: EXACT MATCH with Sarvesh's real tunnel JSON", 1),
        ("All 4 fields (source, destination, SDP ID, name) correct", 1),
        "",
        "Test 2: ePIPE Service",
        ('Input: "Create E-Pipe for customer 10, 192.168.0.37 to 192.168.0.16, VLAN 1001"', 1),
        ("Result: EXACT MATCH with Sarvesh's real epipe JSON", 1),
        ("All 16 fields correct, including SDP bidirectionality", 1),
        "",
        "Both outputs are directly deployable to NSP RESTConf API",
    ])

    # === Slide 15: Supported Intent Types ===
    add_table_slide(prs, "Supported Intent Types",
        ["Intent Type", "Service", "Complexity", "Status"],
        [
            ["Tunnel", "MPLS SDP with BGP", "~6 parameters", "Validated"],
            ["EPIPE", "Point-to-Point Ethernet (E-Line)", "~16 parameters", "Validated"],
            ["VPRN (1 site)", "L3 VPN, single site", "~15 parameters", "Demonstrated"],
            ["VPRN (2+ sites)", "L3 VPN, multi-site", "~30+ parameters", "Demonstrated"],
        ]
    )

    # === Slide 16: Conclusion ===
    add_title_slide(
        prs,
        "Summary",
        "Natural Language -> NSP Intent JSON\n"
        "Fine-tuned Qwen3.5-9B | BF16 + LoRA | 52 min training\n"
        "Validated against real operational data | Ready for deployment"
    )

    # Save
    output_path = "/home/nextron/nsp_intent_ft/NSP_Intent_FT_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")


if __name__ == "__main__":
    main()
